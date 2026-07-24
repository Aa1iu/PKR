"""
文档入库脚本（Phase 1 验证用）

用法：
    python scripts/ingest.py <file_path> --kb-id <kb_id>

将文档解析、分块、embedding 后存入 ChromaDB。
"""

import argparse
import os
import sys

# 将 backend 加入路径
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "backend"))


def parse_document(file_path: str) -> str:
    """解析文档，返回纯文本"""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return text

    elif ext == ".docx":
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs)

    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)
        return "\n".join(texts)

    elif ext in (".txt", ".md"):
        import chardet
        with open(file_path, "rb") as f:
            raw = f.read()
        encoding = chardet.detect(raw)["encoding"] or "utf-8"
        return raw.decode(encoding)

    else:
        raise ValueError(f"不支持的文件类型: {ext}")


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    """将长文本切分为重叠块"""
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", "。", ".", " ", ""],
    )
    return splitter.split_text(text)


def main():
    parser = argparse.ArgumentParser(description="文档入库")
    parser.add_argument("file_path", help="文档路径")
    parser.add_argument("--kb-id", required=True, help="目标知识库 ID")
    parser.add_argument("--chunk-size", type=int, default=500)
    parser.add_argument("--chunk-overlap", type=int, default=50)
    args = parser.parse_args()

    print(f"[1/3] 解析文档: {args.file_path}")
    text = parse_document(args.file_path)
    print(f"  解析完成，共 {len(text)} 字符")

    print(f"[2/3] 文本分块 (chunk_size={args.chunk_size}, overlap={args.chunk_overlap})")
    chunks = chunk_text(text, args.chunk_size, args.chunk_overlap)
    print(f"  分块完成，共 {len(chunks)} 块")

    print("[3/3] Embedding + 入库 ChromaDB")
    # TODO: Phase 1 — ChromaDB 集成
    print(f"  (TODO) 将 {len(chunks)} 个块写入 ChromaDB，kb_id={args.kb_id}")

    print("\n✅ 入库完成（ChromаDB 写入功能待 Phase 1 实现）")


if __name__ == "__main__":
    main()
