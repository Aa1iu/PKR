"""文档解析与分块工具

供 API 路由和 CLI 脚本共享使用。
"""

import os


def parse_document(file_path: str) -> str:
    """解析文档，返回纯文本"""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        import fitz
        doc = fitz.open(file_path)
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return text

    elif ext == ".docx":
        from docx import Document as DocxDocument
        doc = DocxDocument(file_path)
        return "\n".join(p.text for p in doc.paragraphs)

    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)
        return "\n".join(texts)

    elif ext in (".txt", ".md"):
        import chardet
        with open(file_path, "rb") as f:
            raw = f.read()
        encoding = chardet.detect(raw)["encoding"] or "utf-8"
        return raw.decode(encoding)

    else:
        raise ValueError(f"不支持的文件类型: {ext}")


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    """将长文本切分为重叠块"""
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", "。", ".", " ", ""],
    )
    return splitter.split_text(text)


def estimate_page_num(chunk_index: int, chunk_size: int = 500,
                      chars_per_page: int = 800) -> int:
    """粗略估算页码"""
    return max(1, (chunk_index * chunk_size) // chars_per_page + 1)


def jaccard_deduplicate(results: list[dict], threshold: float = 0.8) -> list[dict]:
    """词级 Jaccard 文本去重（中文用 2-gram，英文用空格分词）

    相似度 > threshold 时保留 score 最高的一条。
    比字符级 Jaccard 更准确：不会因中文字符集重叠而误判。
    """
    if not results:
        return results

    def tokenize(text: str) -> set:
        # 中文：2-gram；英文/数字保持词级
        import re
        # 提取中文字符
        chinese = re.findall(r'[一-鿿]', text)
        # 中文 2-gram
        c_bigrams = {chinese[i] + chinese[i + 1] for i in range(len(chinese) - 1)}
        # 英文/数字按空格和标点分词
        english = set(re.findall(r'[a-zA-Z0-9]+', text.lower()))
        return c_bigrams | english

    def jaccard_sim(a: str, b: str) -> float:
        set_a, set_b = tokenize(a), tokenize(b)
        if not set_a or not set_b:
            return 0.0
        return len(set_a & set_b) / len(set_a | set_b)

    kept = []
    for r in results:
        is_dup = False
        for i, k in enumerate(kept):
            if jaccard_sim(r["chunk_text"], k["chunk_text"]) > threshold:
                is_dup = True
                if r["score"] > k["score"]:
                    kept[i] = r
                break
        if not is_dup:
            kept.append(r)
    return kept
